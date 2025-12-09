from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
slide_w, slide_h = prs.slide_width, prs.slide_height

colors = {
    "bg": RGBColor(14, 24, 44),
    "panel": RGBColor(26, 38, 64),
    "accent": RGBColor(95, 173, 86),
    "accent2": RGBColor(255, 194, 102),
    "text": RGBColor(235, 241, 248),
    "muted": RGBColor(186, 200, 220),
}


def add_background(slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = colors["bg"]
    rect.line.fill.background()


def add_header(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["accent"]
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = text
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = colors["bg"]
    p.alignment = PP_ALIGN.CENTER


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), slide_w - Inches(1.6), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = colors["text"]
    p.line_spacing = 1.05
    if subtitle:
        para = tf.add_paragraph()
        para.text = subtitle
        para.font.size = Pt(22)
        para.font.color.rgb = colors["muted"]
        para.space_before = Pt(10)


def add_bullets(slide, heading, items, left=Inches(0.8), top=Inches(1.0)):
    box = slide.shapes.add_textbox(left, top, slide_w - Inches(1.6), slide_h - top - Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = heading
    p0 = tf.paragraphs[0]
    p0.font.size = Pt(30)
    p0.font.bold = True
    p0.font.color.rgb = colors["text"]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = colors["muted"]
        p.level = 1
        p.line_spacing = 1.05
    return box


def add_panel_text(slide, text, left, top, width, height, title=None):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = colors["panel"]
    panel.line.color.rgb = colors["accent"]
    tf = panel.text_frame
    tf.word_wrap = True
    if title:
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = colors["text"]
        p.space_after = Pt(6)
    else:
        tf.text = ""
    for line in text.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(19)
        p.font.color.rgb = colors["muted"]
        p.level = 1
    return panel


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# Slide 1 title
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s1)
add_title(s1, "Week 16: Logic Gates \u2192 Python Conditionals", "Slow, step-by-step mapping from 0/1 to if/elif/else")
set_notes(s1, "Welcome back. Today we bridge last week\u2019s logic gates to Python conditionals. We will move slowly, mapping 0 and 1 to False and True, then building if, elif, and else step by step.")

# Slide 2 objectives
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s2)
add_header(s2, "Objectives")
add_bullets(s2, "By the end, you can:", [
    "Translate gate-style logic into Python boolean expressions.",
    "Use comparisons and and/or/not to build conditions.",
    "Write clear if/elif/else blocks for mutually exclusive cases.",
    "Explain short-circuiting and truthy/falsy values.",
])
set_notes(s2, "Objectives: translate gate logic to Python, combine comparisons with and or not, write if/elif/else cleanly, and explain short-circuiting plus truthy versus falsy.")

# Slide 3 bridge
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s3)
add_header(s3, "Bridge: Gates to Code")
add_panel_text(s3, "A NOT gate flips a 0 to 1; Python uses not to flip False to True.\nAND gate needs both inputs = 1; Python and needs both True.\nOR gate needs any input = 1; Python or needs any True.\nXOR is exactly one input different; Python can use != on booleans.", Inches(0.8), Inches(1.0), slide_w - Inches(1.6), Inches(3.6), title="Key mappings")
set_notes(s3, "Remind them: NOT maps to not, AND to and, OR to or, XOR to != on booleans. Gates output 0 or 1; Python uses False and True. Same patterns, new syntax.")

# Slide 4 booleans and comparisons
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s4)
add_header(s4, "Booleans & Comparisons")
add_bullets(s4, "Common comparisons:", ["== (equal), != (not equal)", ">, >=, <, <= for numbers", "Combine with and/or/not"], left=Inches(0.8), top=Inches(1.0))
add_panel_text(s4, "Example: score = 87\nscore >= 70 \u2192 True\nscore == 100 \u2192 False\nnot(score < 60) \u2192 True", Inches(0.8), Inches(3.0), slide_w - Inches(1.6), Inches(2.2), title="Try it")
set_notes(s4, "Show comparisons: equal, not equal, greater/less. Combine with and or not. Walk through score equals 87 examples: score >= 70 is True, score == 100 is False, not(score < 60) is True.")

# Slide 5 logical operators
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s5)
add_header(s5, "Logical Operators")
add_bullets(s5, "Remember:", ["and \u2192 both conditions must be True", "or \u2192 any condition True passes", "not \u2192 flips the boolean", "Use parentheses to make intent clear"], left=Inches(0.8), top=Inches(1.0))
add_panel_text(s5, "2FA unlocks when password_ok and code_ok.\nAlerts trigger on motion or window sensor.\nMuted if not notifications_enabled.", Inches(0.8), Inches(3.0), slide_w - Inches(1.6), Inches(2.2), title="Gate-style examples")
set_notes(s5, "Link to real cases: 2FA uses and. Alarms use or. Muting uses not. Encourage parentheses for clarity.")

# Slide 6 short-circuit
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s6)
add_header(s6, "Short-Circuiting")
add_panel_text(s6, "and stops early if the first part is False (like an AND gate seeing 0).\nor stops early if the first part is True (like an OR gate seeing 1).\nUse this to avoid calling code that might be slow or error-prone.", Inches(0.8), Inches(1.0), slide_w - Inches(1.6), Inches(2.8), title="Why it matters")
add_panel_text(s6, "Example:\nuser_logged_in and charge_card()\nIf user_logged_in is False, charge_card never runs.", Inches(0.8), Inches(3.3), slide_w - Inches(1.6), Inches(1.8))
set_notes(s6, "Explain short-circuit. If the left side of and is False, Python skips the right side. Same for or: if left is True, skip right. Use it to avoid errors like calling charge_card when not logged in.")

# Slide 7 truthy falsy
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s7)
add_header(s7, "Truthy vs Falsy")
add_bullets(s7, "Falsy values in Python:", ["0, 0.0", "Empty string '', empty list [], empty dict {}", "None"], left=Inches(0.8), top=Inches(1.0))
add_panel_text(s7, "Everything else is truthy.\nWrap with bool(value) to see how Python treats it.\nExample: if cart_items: print('Checkout!')", Inches(0.8), Inches(2.8), slide_w - Inches(1.6), Inches(2.4))
set_notes(s7, "Review falsy: zero, empty string, empty list, empty dict, None. Everything else is truthy. Show bool(value) and how if cart_items uses this.")

# Slide 8 if else structure
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s8)
add_header(s8, "if / else pattern")
add_panel_text(s8, "if condition:\n    do_something()\nelse:\n    do_other()\n\nIndentation shows what belongs to the block. Start simple, then add elif when you have 3+ cases.", Inches(0.8), Inches(1.0), slide_w - Inches(1.6), Inches(3.4))
set_notes(s8, "Show the basic shape. Indentation matters. Use else for the opposite branch. Keep it simple first, then add elif when you need more than two paths.")

# Slide 9 elif
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s9)
add_header(s9, "elif for multiple paths")
add_panel_text(s9, "if score >= 90: grade = 'A'\nelif score >= 80: grade = 'B'\nelif score >= 70: grade = 'C'\nelse: grade = 'Needs work'\n\nFirst true condition wins; later ones are skipped.", Inches(0.8), Inches(1.0), slide_w - Inches(1.6), Inches(3.4))
set_notes(s9, "Walk through the grading example. Emphasize that the first true condition wins, so order matters. Tie back to mutually exclusive cases.")

# Slide 10 XOR toggle
s10 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s10)
add_header(s10, "XOR-style Toggle")
add_panel_text(s10, "Light toggles when button is pressed:\nif button_pressed:\n    light_on = not light_on\n\nThis mirrors XOR: output changes when exactly one input differs.", Inches(0.8), Inches(1.0), slide_w - Inches(1.6), Inches(2.6))
set_notes(s10, "Show how toggles mimic XOR. When button_pressed is True, flip the light state. If not pressed, keep it. Links gate thinking to a real UI action.")

# Slide 11 common mistakes
s11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s11)
add_header(s11, "Common Mistakes")
add_bullets(s11, "Watch out for:", ["Using multiple if instead of elif (cases overlap)", "Forgetting else \u2192 missing default path", "Mixing == with = (assignment)", "Hard-coding magic numbers; prefer variables"], left=Inches(0.8), top=Inches(1.0))
set_notes(s11, "Warn about using separate if instead of elif, forgetting else, using == versus =, and magic numbers. Encourage clear variable names and comments.")

# Slide 12 practice prompts
s12 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s12)
add_header(s12, "Practice Prompts")
add_bullets(s12, "Try in notebook:", ["Write a login check: locked? else 2FA? else deny.", "Alert rule: motion or window, but NOT if system_muted.", "Shipping tiers: free over $50; else $5 under $20; else $2."], left=Inches(0.8), top=Inches(1.0))
set_notes(s12, "Point to the notebook exercises and invite them to try these prompts live. Keep it slow: predict, then run.")

# Slide 13 notebook how-to
s13 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s13)
add_header(s13, "Using the Colab Notebook")
add_bullets(s13, "Steps:", ["Open week_16/Week16_Python_Conditionals.ipynb in Colab.", "Run cells top to bottom; fill TODOs and add prints if helpful.", "Keep comments explaining why your condition is correct."], left=Inches(0.8), top=Inches(1.0))
set_notes(s13, "Explain where the notebook is and to run top to bottom. Encourage adding prints and comments while they work through TODOs.")

# Slide 14 deliverable
s14 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s14)
add_header(s14, "Deliverable")
add_bullets(s14, "Due next class:", ["Complete the Colab notebook exercises and run all cells.", "Keep code commented and include short written answers where prompted.", "Export notebook or PDF and upload to LMS."], left=Inches(0.8), top=Inches(1.0))
set_notes(s14, "Tell them the deliverable: finish the notebook, keep comments, include written answers, export and upload before next class.")

prs.save("week_16/Week16_Conditionals.pptx")
print("Created week_16/Week16_Conditionals.pptx with", len(prs.slides), "slides")
