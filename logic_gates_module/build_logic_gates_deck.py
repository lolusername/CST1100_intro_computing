from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

# Color palette
colors = {
    "bg": RGBColor(12, 22, 38),
    "panel": RGBColor(25, 41, 63),
    "accent": RGBColor(18, 146, 166),
    "accent2": RGBColor(255, 183, 77),
    "text": RGBColor(237, 243, 252),
    "muted": RGBColor(184, 201, 226),
}

prs = Presentation()
slide_w, slide_h = prs.slide_width, prs.slide_height


def add_background(slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    rect.fill.solid(); rect.fill.fore_color.rgb = colors["bg"]
    rect.line.fill.background()


def add_header(slide, label):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = colors["accent"]
    bar.line.fill.background()
    tf = bar.text_frame; tf.text = label; tf.word_wrap = True
    p = tf.paragraphs[0]; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = colors["bg"]
    p.alignment = PP_ALIGN.CENTER


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), slide_w - Inches(1.8), Inches(2.5))
    tf = box.text_frame; tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = colors["text"]; p.line_spacing = 1.05
    if subtitle:
        para = tf.add_paragraph(); para.text = subtitle; para.font.size = Pt(22); para.font.color.rgb = colors["muted"]; para.space_before = Pt(12)


def add_panel(slide, left, top, width, height):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid(); panel.fill.fore_color.rgb = colors["panel"]
    panel.line.color.rgb = colors["accent"]; panel.line.width = Pt(1.2)
    return panel


def add_bullets(slide, title, bullets, left=Inches(0.8), top=Inches(1.0)):
    box = slide.shapes.add_textbox(left, top, slide_w - Inches(1.6), slide_h - top - Inches(0.7))
    tf = box.text_frame; tf.word_wrap = True
    tf.text = title
    p0 = tf.paragraphs[0]; p0.font.size = Pt(30); p0.font.bold = True; p0.font.color.rgb = colors["text"]
    for item in bullets:
        p = tf.add_paragraph(); p.text = item; p.font.size = Pt(20); p.font.color.rgb = colors["muted"]; p.level = 1; p.line_spacing = 1.1
    return box


def add_truth_table(slide, headers, rows, left, top, width, height):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = colors["accent"]
        para = cell.text_frame.paragraphs[0]; para.font.size = Pt(17); para.font.bold = True; para.font.color.rgb = colors["bg"]; para.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c); cell.text = str(value); cell.fill.solid(); cell.fill.fore_color.rgb = colors["panel"]
            para = cell.text_frame.paragraphs[0]; para.font.size = Pt(17); para.font.color.rgb = colors["text"]; para.alignment = PP_ALIGN.CENTER
    table.first_row = True
    return table


def add_gate_diagram(slide, gate_label, left, top, input_labels=("A", "B"), output_label="Output"):
    # Inputs
    for idx, name in enumerate(input_labels):
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.9) * idx, Inches(0.65), Inches(0.65))
        circ.fill.solid(); circ.fill.fore_color.rgb = colors["panel"]; circ.line.color.rgb = colors["muted"]
        tf = circ.text_frame; tf.text = name; tf.word_wrap = True
        p = tf.paragraphs[0]; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = colors["text"]; p.alignment = PP_ALIGN.CENTER
    gate = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, left + Inches(1.05), top + Inches(0.25), Inches(1.6), Inches(1.2))
    gate.fill.solid(); gate.fill.fore_color.rgb = colors["accent"]; gate.line.color.rgb = colors["accent2"]
    tf = gate.text_frame; tf.text = gate_label; tf.word_wrap = True
    p = tf.paragraphs[0]; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = colors["bg"]; p.alignment = PP_ALIGN.CENTER
    out = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(3.15), top + Inches(0.6), Inches(0.65), Inches(0.65))
    out.fill.solid(); out.fill.fore_color.rgb = colors["panel"]; out.line.color.rgb = colors["muted"]
    tf = out.text_frame; tf.text = output_label; tf.word_wrap = True
    p = tf.paragraphs[0]; p.font.size = Pt(14); p.font.color.rgb = colors["text"]; p.alignment = PP_ALIGN.CENTER
    # connectors
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(0.65), top + Inches(0.35), left + Inches(1.05), top + Inches(0.55)).line.color.rgb = colors["muted"]
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(0.65), top + Inches(1.25), left + Inches(1.05), top + Inches(1.05)).line.color.rgb = colors["muted"]
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(2.65), top + Inches(0.85), left + Inches(3.15), top + Inches(0.9)).line.color.rgb = colors["muted"]


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# Slide 1: Title
s1 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s1)
add_title(s1, "Logic Gates: From Signals to Systems", "CST1100 – Intro to Computing")
set_notes(s1, "Welcome. Today we unpack logic gates—the small decision-makers inside every computer. We will start with binary signals, cover the core gates and their truth tables, combine them into adders, and connect the same logic to real applications. Please jump in with questions; the aim is confidence, not memorization.")

# Slide 2: Outcomes & Agenda
s2 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s2); add_header(s2, "Outcomes & Agenda")
add_bullets(s2, "By the end, you will:", [
    "Build truth tables for NOT, AND, OR, NAND, NOR, XOR, XNOR.",
    "Sketch gates and trace a small circuit end-to-end.",
    "Connect gate logic to auth, alerts, and basic arithmetic.",
    "Debug by comparing expected tables to observed outputs.",
])
set_notes(s2, "Roadmap: quick binary refresher; core gates and truth tables; combine gates into adders; then real-world examples and your reflection assignment.")

# Slide 3: Binary Signals
s3 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s3); add_header(s3, "Binary Signals")
panel = add_panel(s3, Inches(0.7), Inches(1.0), slide_w - Inches(1.4), Inches(4.9))
box = s3.shapes.add_textbox(Inches(1.0), Inches(1.3), slide_w - Inches(2.0), Inches(4.3))
tf = box.text_frame; tf.word_wrap = True
lines = [
    "Digital circuits use two stable voltage ranges.",
    "0 (LOW) ≈ 0 volts → treated as false.",
    "1 (HIGH) ≈ supply voltage → treated as true.",
    "Gates read inputs and instantly output a new 0/1.",
    "Truth tables list every input combo with its output.",
]
for i, line in enumerate(lines):
    p = tf.add_paragraph() if i else tf.paragraphs[0]
    p.text = line; p.font.size = Pt(22 if i == 0 else 20); p.font.color.rgb = colors["muted"] if i else colors["text"]
set_notes(s3, "Digital electronics stay reliable by using two ranges: low is near zero volts, high is near the supply voltage. Gates read those levels and output a new level. Truth tables list every input combination with the output that follows so we can reason about a gate quickly.")

# Slide 4: NOT
s4 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s4); add_header(s4, "NOT Gate")
add_truth_table(s4, ["Input", "Output"], [[0, 1], [1, 0]], Inches(0.7), Inches(1.0), Inches(2.7), Inches(1.3))
add_gate_diagram(s4, "NOT", Inches(4.6), Inches(1.0), input_labels=("Input", " "), output_label="Flips")
box = s4.shapes.add_textbox(Inches(0.7), Inches(2.6), slide_w - Inches(1.4), Inches(2.8))
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Behavior"; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = colors["text"]
p2 = tf.add_paragraph(); p2.text = "Outputs the opposite of the input."; p2.font.size = Pt(20); p2.font.color.rgb = colors["muted"]; p2.level = 1
p3 = tf.add_paragraph(); p3.text = "Useful for NOT expired, NOT muted, NOT blocked checks."; p3.font.size = Pt(20); p3.font.color.rgb = colors["muted"]; p3.level = 1
set_notes(s4, "The NOT gate is an inverter. A zero in becomes a one out, and a one in becomes a zero out. We use it whenever we need to express 'not expired', 'not muted', or 'not blocked'.")

# Slide 5: AND
s5 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s5); add_header(s5, "AND Gate")
add_truth_table(s5, ["A", "B", "Out"], [[0,0,0],[0,1,0],[1,0,0],[1,1,1]], Inches(0.7), Inches(1.0), Inches(3.4), Inches(1.4))
add_gate_diagram(s5, "AND", Inches(4.8), Inches(1.0))
add_bullets(s5, "Common uses", [
    "Two-factor auth: password AND code.",
    "Safety interlock: door closed AND sensor aligned.",
    "Rate limit: under quota AND in time window.",
], left=Inches(0.7), top=Inches(2.8))
set_notes(s5, "An AND gate outputs one only when every input is one. Think of two switches in series—both must be on for current to flow. In software, that is a password AND a code, or a safety system needing multiple sensors aligned.")

# Slide 6: OR
s6 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s6); add_header(s6, "OR Gate")
add_truth_table(s6, ["A", "B", "Out"], [[0,0,0],[0,1,1],[1,0,1],[1,1,1]], Inches(0.7), Inches(1.0), Inches(3.4), Inches(1.4))
add_gate_diagram(s6, "OR", Inches(4.8), Inches(1.0))
add_bullets(s6, "Common uses", [
    "Alert: motion OR window sensor triggers alarm.",
    "Routing: email OR SMS OR push sends a notice.",
    "Search: tag A OR tag B matches item.",
], left=Inches(0.7), top=Inches(2.8))
set_notes(s6, "An OR gate outputs one when any input is one. Picture switches in parallel—if either conducts, the lamp lights. We rely on OR for alarms, notifications, or any place where multiple signals can trigger the same response.")

# Slide 7: NAND & NOR
s7 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s7); add_header(s7, "NAND and NOR")
add_truth_table(s7, ["A", "B", "NAND"], [[0,0,1],[0,1,1],[1,0,1],[1,1,0]], Inches(0.7), Inches(1.0), Inches(3.1), Inches(1.4))
add_truth_table(s7, ["A", "B", "NOR"], [[0,0,1],[0,1,0],[1,0,0],[1,1,0]], Inches(4.2), Inches(1.0), Inches(3.1), Inches(1.4))
add_bullets(s7, "Why they matter", [
    "Each is universal: you can build any circuit from just NANDs or just NORs.",
    "NAND is cheap and fast in hardware; NOR appears in reset and default-low logic.",
    "Both mirror code patterns like !(A && B) and !(A || B).",
], left=Inches(0.7), top=Inches(2.8))
set_notes(s7, "NAND flips the AND output; NOR flips the OR output. Each one alone is universal—any gate can be made from copies of one of these. Hardware often prefers NAND for cost and speed. NOR shows up where you only want a high signal when everything else is low, like reset lines."
)

# Slide 8: XOR & XNOR
s8 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s8); add_header(s8, "XOR and XNOR")
add_truth_table(s8, ["A", "B", "XOR"], [[0,0,0],[0,1,1],[1,0,1],[1,1,0]], Inches(0.7), Inches(1.0), Inches(3.1), Inches(1.4))
add_truth_table(s8, ["A", "B", "XNOR"], [[0,0,1],[0,1,0],[1,0,0],[1,1,1]], Inches(4.2), Inches(1.0), Inches(3.1), Inches(1.4))
add_bullets(s8, "Where they show up", [
    "XOR: exactly one input true — parity checks and toggles.",
    "XNOR: inputs match — equality checks and comparators.",
    "UI toggles use XOR; storage checks often use XNOR.",
], left=Inches(0.7), top=Inches(2.8))
set_notes(s8, "XOR outputs one only when the inputs differ, so it is great for parity checks and toggles. XNOR outputs one when the inputs match, behaving like an equality check. You see XOR in toggles and parity bits; you see XNOR in comparators and change detection.")

# Slide 9: Adders
s9 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s9); add_header(s9, "Composing Gates: Adders")
add_truth_table(s9, ["A", "B", "Sum", "Carry"], [[0,0,0,0],[0,1,1,0],[1,0,1,0],[1,1,0,1]], Inches(0.6), Inches(1.0), Inches(4.2), Inches(1.6))
box1 = s9.shapes.add_textbox(Inches(5.0), Inches(1.0), slide_w - Inches(5.6), Inches(2.0))
tf1 = box1.text_frame; tf1.word_wrap = True
p = tf1.paragraphs[0]; p.text = "Half adder"; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = colors["text"]
p = tf1.add_paragraph(); p.text = "Sum = XOR(A, B)"; p.font.size = Pt(19); p.font.color.rgb = colors["muted"]; p.level = 1
p = tf1.add_paragraph(); p.text = "Carry = AND(A, B)"; p.font.size = Pt(19); p.font.color.rgb = colors["muted"]; p.level = 1
box2 = s9.shapes.add_textbox(Inches(0.6), Inches(2.9), slide_w - Inches(1.2), Inches(2.6))
tf2 = box2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]; p.text = "Full adder"; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = colors["text"]
for line in [
    "Adds A + B + carry-in.",
    "Sum = XOR(A, B, carry-in); carry-out = 1 when two or more inputs are 1.",
    "Chain full adders to build multi-bit arithmetic units.",
]:
    para = tf2.add_paragraph(); para.text = line; para.font.size = Pt(19); para.font.color.rgb = colors["muted"]; para.level = 1
set_notes(s9, "Combining gates gives arithmetic. A half adder uses XOR for the sum bit and AND for the carry. A full adder adds a carry-in; if two or more inputs are one, the carry-out is one. Chain full adders to make the arithmetic logic in a CPU.")

# Slide 10: Two-Factor Login
s10 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s10); add_header(s10, "Diagram: Two-Factor Login")
add_gate_diagram(s10, "AND", Inches(0.9), Inches(1.0), input_labels=("Password OK", "Code OK"), output_label="Unlock")
callout = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(1.5), Inches(4.8), Inches(1.6))
callout.fill.solid(); callout.fill.fore_color.rgb = colors["panel"]; callout.line.color.rgb = colors["accent"]
tf = callout.text_frame; tf.word_wrap = True; tf.text = "Output: unlock when password AND code are both true."; tf.paragraphs[0].font.size = Pt(20); tf.paragraphs[0].font.color.rgb = colors["text"]; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
set_notes(s10, "Two-factor login is an AND gate. Input A is password correct. Input B is verification code correct. Only when both are true does the output unlock. If either fails, output stays zero. Same pattern as an if-statement requiring two conditions.")

# Slide 11: Real-World Patterns
s11 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s11); add_header(s11, "Real-World Patterns")
add_bullets(s11, "Examples:", [
    "Safety: guards closed AND emergency stop NOT pressed.",
    "Alerts: smoke OR heat trips alarm; sprinkler adds AND overheat to cut false alarms.",
    "Networking: allow when source trusted AND destination allowed AND NOT throttled.",
    "UI toggle: dark mode XOR switch flips each click.",
])
set_notes(s11, "Gate logic powers safety systems, alerts, networking, and UI toggles. A factory press needs guards closed AND the emergency stop NOT pressed. Alarms may use OR to trigger, with an AND on temperature to cut false alarms. Firewalls combine AND, OR, and NOT to decide packet flow. A toggle button is an XOR—it flips state each click.")

# Slide 12: Inside the CPU
s12 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s12); add_header(s12, "Inside the CPU")
add_bullets(s12, "Gate-driven blocks:", [
    "Instruction decode: AND masks on opcode bits select a path.",
    "Control: OR chains raise flags like write-enable or branch.",
    "ALU: chains of adders implement +, -, and comparisons.",
    "Caches: tag match uses XNOR/XOR for equality and difference.",
])
set_notes(s12, "A CPU is a sea of gates. Instruction decoders AND opcode bits with masks to route control signals. OR chains raise flags like write-enable or branch. The ALU is built from the adders we saw. Caches use XNOR or XOR to compare addresses quickly.")

# Slide 13: Bitwise in Code
s13 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s13); add_header(s13, "Bitwise Operators in Code")
box = s13.shapes.add_textbox(Inches(0.8), Inches(1.0), slide_w - Inches(1.6), Inches(4.8))
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Example (Python)"; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = colors["text"]
code = (
    "mask = 0b1010  # feature flags\n"
    "user = 0b1110  # current settings\n"
    "# AND: is flag 0b0010 on?\n"
    "enabled = bool(user & 0b0010)\n"
    "# XOR: toggle a flag\n"
    "user ^= 0b1000\n"
    "# OR: force-enable a flag\n"
    "user |= 0b0001\n"
)
p = tf.add_paragraph(); p.text = code; p.font.name = "Courier New"; p.font.size = Pt(18); p.font.color.rgb = colors["muted"]; p.level = 1
set_notes(s13, "Software mirrors gate logic with bitwise operators. AND checks whether a specific flag is on. XOR toggles a bit: if it was one, it becomes zero; if zero, it becomes one. OR forces a bit on. Understanding gates makes these operators clearer.")

# Slide 14: Debugging
s14 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s14); add_header(s14, "Debugging and Pitfalls")
add_bullets(s14, "Watch for:", [
    "Mixing up XOR and OR — XOR is exactly one.",
    "Forgetting a NOT, so outputs flip.",
    "Truth table typos — one row can change the gate type.",
    "Unclear signal names; label inputs and outputs.",
])
set_notes(s14, "Common mistakes: mixing up XOR and OR—remember XOR means exactly one. Forgetting a NOT is another frequent error. A single wrong row in a truth table can misidentify a gate. Clear signal names help everyone debug faster.")

# Slide 15: Reflection
s15 = prs.slides.add_slide(prs.slide_layouts[6]); add_background(s15); add_header(s15, "Reflection Assignment")
add_bullets(s15, "Submit 200–300 words (see handout):", [
    "Pick one scenario (e.g., 2FA, alarm, toggle, traffic light).",
    "Draw or describe gates; label inputs and output.",
    "Make a 4-row truth table and explain one row plainly.",
    "Say why correctness matters and share one learning/question.",
])
set_notes(s15, "Your reflection connects the theory to something you use. Choose one scenario, map its gates, and include a small truth table. Explain one row in plain language and why correctness matters. Keep it to 200 to 300 words and submit before next class.")

prs.save("logic_gates_module/Logic_Gates_Deep_Dive.pptx")
print("Deck rebuilt: logic_gates_module/Logic_Gates_Deep_Dive.pptx with", len(prs.slides), "slides")
